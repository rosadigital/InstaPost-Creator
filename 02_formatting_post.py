import os
from pptx import Presentation
from pptx.util import Cm
from PIL import Image

class Main():
    def creatingpost(self, text_created):
        print('--------------------------------')
        print('Starting to format your post.')

        current_directory = os.path.dirname(os.path.realpath(__file__))

        '''creating pptx'''
        prs = Presentation(current_directory + '\\modelforedition.pptx')
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)

        '''adding image'''
        image = Image.open(text_created + '.jpeg', 'r')
        image_size = image.size
        width = image_size[0]
        width_in_cm = width * 0.0264583  # converting pixel to cm
        height = image_size[1]
        height_in_cm = height * 0.0264583  # converting pixel to cm
        img_path = text_created + '.jpeg'
        left = Cm((19.05 - width_in_cm) / 2)
        top = Cm((13.33 - height_in_cm) / 2)
        pic = slide.shapes.add_picture(img_path, left, top)

        '''writing title'''
        title = slide.shapes.title
        title.text = text_created

        '''writing subtitle'''
        subtitle = slide.placeholders[1]
        subtitle_for_post = open(current_directory + "\\" + text_created + "_for_post.txt", "r", encoding="utf-8").read()
        subtitle.text = subtitle_for_post

        '''saving pptx'''
        prs.save(text_created + '.pptx')

        print('\nUhuuuu!!! Process concluded.')


if __name__ =='__main__':
    text_created = input("Write the theme that you wrote while running 01_creating_text.py."
                         "\n I.e if when running 01_creating_text.py you wrote Christina_Aguilera, now write againg Christina_Aguilera."
                         "\n Write here, and press enter >>>: ")
    start = Main()
    start.creatingpost(text_created)